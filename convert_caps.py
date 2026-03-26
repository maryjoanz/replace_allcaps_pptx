#!/usr/bin/env python3
"""
convert_caps.py
───────────────
Finds text in a PowerPoint (.pptx) that uses All Caps or Small Caps font
formatting and converts it to mixed (title) case, removing the caps formatting.

Usage:
    python convert_caps.py input.pptx
    python convert_caps.py *.pptx --keep AI API NASA
    python convert_caps.py *.pptx --outdir converted/

Single file:   saves to <input>_mixed_case.pptx (or specify -o name.pptx)
Batch (2+ files): saves each to <input>_mixed_case.pptx, or into --outdir

The --keep flag adds extra words that should keep their casing.  A built-in
list of common acronyms is always applied (AI, API, CEO, etc.).  Edit the
KEEP_WORDS set below to permanently add or remove entries.

Requirements:
    pip install python-pptx
"""

import re
import sys
import glob
import argparse
from pathlib import Path
from lxml import etree
from pptx import Presentation

# ── Acronyms / words to preserve ────────────────────────────────
# These words will keep their exact casing after title-case conversion.
# Add or remove entries as needed.  Case-sensitive: "AI" means the
# uppercase form "AI" is restored whenever "Ai" appears after title-casing.
KEEP_WORDS = {
    "AI", "API", "APIs", "AWS", "CEO", "CFO", "CIO", "CTO", "COO",
    "CSV", "DB", "DevOps", "DNS", "EU", "FAQ", "GPU", "GPUs",
    "HR", "HTML", "HTTP", "HTTPS", "I", "ID", "IDs", "iOS", "IoT",
    "IP", "IT", "JSON", "KPI", "KPIs", "LLC", "MBA", "ML", "MVP",
    "NASA", "NDA", "OKR", "OKRs", "OS", "PDF", "PhD", "PR", "QA",
    "ROI", "SaaS", "SDK", "SEO", "SQL", "SSD", "SSO", "UI", "UK",
    "URL", "URLs", "US", "USA", "USB", "UX", "VPN", "XML",
}

# PowerPoint XML namespace for drawingML
_nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def _get_cap_attr(run_element):
    """Return the 'cap' attribute value from a run's <a:rPr> element, or None."""
    rPr = run_element.find("a:rPr", _nsmap)
    if rPr is not None:
        return rPr.get("cap")
    return None


def _remove_cap_attr(run_element):
    """Remove the 'cap' attribute from a run's <a:rPr> element."""
    rPr = run_element.find("a:rPr", _nsmap)
    if rPr is not None and "cap" in rPr.attrib:
        del rPr.attrib["cap"]


def _to_mixed_case(text, keep_words):
    """
    Convert text to title case, then:
      1. Fix post-apostrophe capitalisation (today's, don't, etc.).
      2. Restore any words listed in *keep_words* to their exact casing.
    """
    titled = text.title()
    # Fix apostrophe issue
    titled = re.sub(r"(?<=[\u2019'])\w", lambda m: m.group().lower(), titled)

    # Build a lookup: lowered form -> desired form
    if keep_words:
        lookup = {w.lower(): w for w in keep_words}
        titled = re.sub(
            r"\b[A-Za-z]+\b",
            lambda m: lookup.get(m.group().lower(), m.group()),
            titled,
        )
    return titled


def _iter_all_text_frames(prs):
    """Yield every text frame in the presentation (slides, layouts, masters, notes)."""
    # Slide content
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                yield shape.text_frame
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        yield cell.text_frame
        # Slide notes
        if slide.has_notes_slide:
            for shape in slide.notes_slide.shapes:
                if shape.has_text_frame:
                    yield shape.text_frame

    # Slide layouts & masters (optional – uncomment if you want those too)
    # for layout in prs.slide_layouts:
    #     for shape in layout.placeholders:
    #         if shape.has_text_frame:
    #             yield shape.text_frame


def convert_caps_to_mixed(input_path, output_path=None, keep_words=None):
    """
    Open a .pptx, find any run with All Caps ('all') or Small Caps ('small')
    formatting, convert its text to title case, and remove the caps attribute.
    Returns the number of runs that were converted.
    """
    if keep_words is None:
        keep_words = KEEP_WORDS
    else:
        keep_words = KEEP_WORDS | set(keep_words)

    input_path = Path(input_path)
    if output_path is None:
        output_path = input_path.with_stem(input_path.stem + "_mixed_case")
    else:
        output_path = Path(output_path)

    prs = Presentation(str(input_path))
    converted_count = 0

    for text_frame in _iter_all_text_frames(prs):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                cap_value = _get_cap_attr(run._r)
                if cap_value in ("all", "small"):
                    original = run.text
                    run.text = _to_mixed_case(original, keep_words)
                    _remove_cap_attr(run._r)
                    converted_count += 1
                    print(
                        f"  [{cap_value.upper()} CAPS] "
                        f'"{original}" -> "{run.text}"'
                    )

    prs.save(str(output_path))
    return converted_count, output_path


def main():
    parser = argparse.ArgumentParser(
        description="Convert All Caps / Small Caps formatted text in "
                    "PowerPoint files to mixed case."
    )
    parser.add_argument("inputs", nargs="+", metavar="FILE",
                        help="One or more .pptx files (supports wildcards)")
    parser.add_argument("-o", "--output", default=None,
                        help="Output filename (single-file mode only)")
    parser.add_argument("--outdir", default=None,
                        help="Output directory for batch mode (created if needed)")
    parser.add_argument("--keep", nargs="+", metavar="WORD", default=[],
                        help="Extra words to preserve casing for "
                             "(e.g. --keep ACME NLP GenAI)")
    args = parser.parse_args()

    # Expand wildcards (Windows doesn't do this automatically)
    expanded = []
    for pattern in args.inputs:
        matches = glob.glob(pattern)
        if matches:
            expanded.extend(matches)
        else:
            expanded.append(pattern)  # keep as-is so the "not found" warning fires
    args.inputs = expanded

    # Validate: -o only makes sense with a single file
    if args.output and len(args.inputs) > 1:
        parser.error("-o/--output can only be used with a single input file. "
                     "Use --outdir for batch mode.")

    # Create outdir if requested
    if args.outdir:
        outdir = Path(args.outdir)
        outdir.mkdir(parents=True, exist_ok=True)

    extra = args.keep if args.keep else None
    total_files = 0
    total_runs = 0

    for input_file in args.inputs:
        input_path = Path(input_file)
        if not input_path.exists():
            print(f"WARNING: {input_file} not found, skipping.")
            continue
        if not input_path.suffix.lower() == ".pptx":
            print(f"WARNING: {input_file} is not a .pptx file, skipping.")
            continue

        # Determine output path
        if args.output:
            out_path = Path(args.output)
        elif args.outdir:
            out_path = Path(args.outdir) / input_path.name
        else:
            out_path = input_path.with_stem(input_path.stem + "_mixed_case")

        print(f"\n{'─' * 60}")
        print(f"Processing: {input_file}")
        count, out = convert_caps_to_mixed(input_file, out_path, extra)
        total_files += 1
        total_runs += count

        if count == 0:
            print(f"  No All Caps or Small Caps runs found.")
        else:
            print(f"  Converted {count} run(s).")
        print(f"  Saved to: {out}")

    # Summary for batch runs
    print(f"\n{'═' * 60}")
    print(f"Done. {total_files} file(s) processed, {total_runs} run(s) converted.")


if __name__ == "__main__":
    main()
