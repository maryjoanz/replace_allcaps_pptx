"""
replace_allcaps_pptx.py

Replaces ALL CAPS text in a PowerPoint file with Mixed Case,
improving accessibility and readability.

Usage:
    python replace_allcaps_pptx.py input.pptx output.pptx [--mode title|sentence]

Requirements:
    pip install python-pptx
"""

import re
import copy
import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt


# ── Acronyms / abbreviations to KEEP uppercase ──────────────────────
PRESERVE_UPPERCASE = {
    "WCAG", "ADA", "PDF", "HTML", "CSS", "URL", "API", "NASA",
    "FBI", "USA", "UK", "EU", "IT", "HR", "CEO", "CFO", "CTO",
    "FAQ", "ID", "PIN", "ISO", "ARIA", "WAI", "VPAT", "JAWS",
    "NVDA", "AI", "ML", "SQL", "JSON", "XML", "HTTP", "HTTPS",
    # Add your own acronyms here
}


def is_all_caps(text: str) -> bool:
    """Return True if the text contains letters and they are ALL uppercase."""
    alpha_chars = [c for c in text if c.isalpha()]
    return len(alpha_chars) >= 2 and all(c.isupper() for c in alpha_chars)


def convert_word(word: str, mode: str = "title") -> str:
    """
    Convert a single ALL-CAPS word to mixed case,
    preserving acronyms and leading/trailing punctuation.
    """
    # Strip punctuation to check the core word
    match = re.match(r'^([^A-Za-z]*)([A-Za-z]+)([^A-Za-z]*)$', word)
    if not match:
        return word  # No alpha characters, return as-is

    prefix, core, suffix = match.groups()

    # Preserve known acronyms
    if core.upper() in PRESERVE_UPPERCASE:
        return word  # Keep it unchanged

    # Only convert if core is ALL CAPS (2+ letters)
    if len(core) >= 2 and core.isupper():
        if mode == "title":
            core = core.capitalize()   # "HELLO" → "Hello"
        elif mode == "sentence":
            core = core.lower()        # Will be re-capitalized at sentence level
    return prefix + core + suffix


def convert_text(text: str, mode: str = "title") -> str:
    """
    Convert ALL CAPS text in a string to mixed case.

    Modes:
        'title'    – Each ALLCAPS word → Capitalized  ("HELLO WORLD" → "Hello World")
        'sentence' – First word capitalized, rest lower ("HELLO WORLD" → "Hello world")
    """
    if not text or not is_all_caps(text):
        return text

    words = text.split(" ")
    converted = [convert_word(w, mode) for w in words]

    result = " ".join(converted)

    # For sentence mode, capitalize the first alpha character
    if mode == "sentence":
        result = _capitalize_first(result)

    return result


def _capitalize_first(text: str) -> str:
    """Capitalize only the first alphabetic character in the string."""
    for i, ch in enumerate(text):
        if ch.isalpha():
            return text[:i] + ch.upper() + text[i + 1:]
    return text


# ── Font-level "All Caps" attribute removal ─────────────────────────
def clear_font_allcaps(font) -> None:
    """
    PowerPoint can apply visual ALL CAPS via the font 'caps' attribute
    (Character Spacing → All Caps). This clears that attribute so the
    actual stored text is what displays.
    """
    # python-pptx doesn't expose `caps` directly, so we edit the XML
    rPr = font._element  # this is the <a:rPr> element
    # The attribute name in OOXML is 'cap' with values 'all', 'small', 'none'
    if rPr is not None and rPr.attrib.get('cap'):
        rPr.attrib.pop('cap', None)


# ── Process a single text frame ─────────────────────────────────────
def process_text_frame(text_frame, mode: str, stats: dict) -> None:
    """Process every paragraph and run inside a text frame."""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            original = run.text

            # 1) Remove font-level "All Caps" formatting
            clear_font_allcaps(run.font._element)

            # 2) Convert the actual stored text
            converted = convert_text(original, mode)

            if converted != original:
                run.text = converted
                stats["runs_changed"] += 1
                stats["details"].append(f'  "{original}" → "{converted}"')

            stats["runs_total"] += 1


# ── Process an entire shape (recursive for groups) ───────────────────
def process_shape(shape, mode: str, stats: dict) -> None:
    """Recursively process shapes, group shapes, and tables."""

    # Grouped shapes
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for child_shape in shape.shapes:
            process_shape(child_shape, mode, stats)
        return

    # Tables
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    process_text_frame(cell.text_frame, mode, stats)
        return

    # Regular text frames
    if shape.has_text_frame:
        process_text_frame(shape.text_frame, mode, stats)


# ── Main processing function ────────────────────────────────────────
def process_presentation(input_path: str, output_path: str, mode: str = "title") -> dict:
    """
    Open a .pptx, replace ALL CAPS text, and save to a new file.

    Returns a stats dictionary with change details.
    """
    prs = Presentation(input_path)

    stats = {
        "slides": 0,
        "runs_total": 0,
        "runs_changed": 0,
        "details": [],
    }

    for slide_index, slide in enumerate(prs.slides, start=1):
        stats["slides"] += 1
        stats["details"].append(f"\n── Slide {slide_index} ──")

        # Process all shapes on the slide
        for shape in slide.shapes:
            process_shape(shape, mode, stats)

        # Process slide notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                process_text_frame(notes_frame, mode, stats)

    # ── Slide masters & layouts (headers/footers, template text) ────
    for master in prs.slide_masters:
        for shape in master.shapes:
            process_shape(shape, mode, stats)
        for layout in master.slide_layouts:
            for shape in layout.shapes:
                process_shape(shape, mode, stats)

    prs.save(output_path)
    return stats


# ── CLI entry point ─────────────────────────────────────────────────
def main():
    parser = argparse.ArgumentParser(
        description="Replace ALL CAPS text in PowerPoint with mixed case."
    )
    parser.add_argument("input", help="Path to input .pptx file")
    parser.add_argument("output", help="Path to output .pptx file")
    parser.add_argument(
        "--mode",
        choices=["title", "sentence"],
        default="title",
        help="Conversion mode: 'title' (default) or 'sentence'",
    )
    parser.add_argument(
        "--verbose", "-v",
        action="store_true",
        help="Print every individual change",
    )
    args = parser.parse_args()

    # Validate input
    if not Path(args.input).exists():
        print(f"Error: '{args.input}' not found.")
        return

    print(f"Processing: {args.input}")
    print(f"Mode:       {args.mode} case")
    print(f"Output:     {args.output}")
    print("─" * 50)

    stats = process_presentation(args.input, args.output, args.mode)

    # ── Report ──────────────────────────────────────────────────────
    if args.verbose:
        for line in stats["details"]:
            print(line)
        print("─" * 50)

    print(f"Slides scanned:  {stats['slides']}")
    print(f"Text runs found: {stats['runs_total']}")
    print(f"Text runs fixed: {stats['runs_changed']}")
    print(f"\n✓ Saved to '{args.output}'")


if __name__ == "__main__":
    main()
